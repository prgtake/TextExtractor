# -*- coding: utf-8 -*-
# =====================================================
#  TextExtractor1 v1.1.0
#  Copyright (c) 2026 Datan (データン)
#  Licensed under the MIT License.
# =====================================================
import tkinter as tk
from tkinter import filedialog, messagebox
import tkinter.simpledialog as sd
import os
import re
import sqlite3
import json
import datetime
import time
import csv
import pandas as pd
import docx
from pptx import Presentation
import extract_msg
from email import policy
from email.parser import BytesParser

from google import genai
from google.genai import types
from google.genai.errors import APIError

# =====================================================
#  設定とAPI初期化
# =====================================================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
CONFIG_FILE = os.path.join(BASE_DIR, "app_config.json")
DEFAULT_MODEL = "gemini-2.5-flash-lite"

# アプリバージョン
APP_VERSION = "1.1.0"

def load_config():
    config = {"GEMINI_API_KEY": os.environ.get("GEMINI_API_KEY"), "MODEL_NAME": DEFAULT_MODEL}
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                config.update(json.load(f))
        except: pass
    return config

def save_config(api_key, model_name):
    try:
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump({"GEMINI_API_KEY": api_key, "MODEL_NAME": model_name}, f)
    except: pass

class TextExtractor:
    def __init__(self, master):
        self.master = master
        master.title(f"TextExtractor1 v{APP_VERSION} ")
        # 修正: 初期高さを 850 から 560 (約2/3) に変更
        master.geometry("650x560")
        master.configure(bg="#f0f2f5")
        self.config = load_config()
        self.api_key = self.config.get("GEMINI_API_KEY")
        
        if not self.api_key or self.api_key == "YOUR_API_KEY_HERE":
            self.api_key = sd.askstring("API キー", "Gemini API キーを入力してください：\n（次回から入力不要になります）", parent=master)
            if not self.api_key:
                print("APIキーが設定されていません。終了します。")
                master.destroy()
                return
            save_config(self.api_key, self.config.get("MODEL_NAME"))

        try:
            self.client = genai.Client(
                api_key=self.api_key,
                http_options={'retry_options': {'attempts': 5, 'initial_delay': 2.0, 'max_delay': 60.0, 'exp_base': 2.0}}
            )
        except Exception as e:
            messagebox.showerror("エラー", f"API初期化失敗: {e}")
            self.client = None

        self.target_folder = None
        self.prompt_text = None
        self.prompt_file_path = None
        self.columns = []

        self.setup_ui()

    def setup_ui(self):
        font_main = ("Meiryo", 10)
        bg_color = "#f0f2f5"
        # 復元: オリジナルのボタンスタイル + 幅を統一
        self.btn_style = {"font": font_main, "padx": 10, "pady": 3, "width": 20}
        
        # 1. フォルダ選択
        f1 = tk.LabelFrame(self.master, text="1. 対象フォルダ選択", bg=bg_color, font=font_main, padx=10, pady=5)
        f1.pack(fill="x", padx=20, pady=5)
        tk.Button(f1, text="フォルダを選択", command=self.select_folder, **self.btn_style).pack(side="left")
        
        # 復元: サブフォルダも含めるチェックボックス
        self.use_recursive = tk.BooleanVar(value=False)
        tk.Checkbutton(f1, text="サブフォルダも含める", variable=self.use_recursive, bg=bg_color, font=font_main).pack(side="left", padx=15)
        
        self.folder_label = tk.Label(f1, text="未選択", fg="blue", bg=bg_color, font=font_main)
        self.folder_label.pack(side="left", padx=10)

        # 2. 指示ファイル選択
        f2 = tk.LabelFrame(self.master, text="2. 指示ファイル選択", bg=bg_color, font=font_main, padx=10, pady=5)
        f2.pack(fill="x", padx=20, pady=5)
        tk.Button(f2, text="指示ファイルを選択", command=self.select_prompt_file, **self.btn_style).pack(side="left")
        tk.Button(f2, text="書き方ヘルプ", command=self.show_help, bg="#ffffcc", **self.btn_style).pack(side="left", padx=5)
        self.prompt_label = tk.Label(f2, text="未選択", fg="blue", bg=bg_color, font=font_main)
        self.prompt_label.pack(side="left", padx=5)

        # 3. モデル・ツール設定
        f3 = tk.LabelFrame(self.master, text="3. 詳細設定", bg=bg_color, font=font_main, padx=10, pady=5)
        f3.pack(fill="x", padx=20, pady=5)
        
        tk.Label(f3, text="使用モデル:", bg=bg_color, font=font_main).grid(row=0, column=0, sticky="w")
        self.model_var = tk.StringVar(value=self.config.get("MODEL_NAME", DEFAULT_MODEL))
        tk.Entry(f3, textvariable=self.model_var, width=25, font=font_main).grid(row=0, column=1, padx=5, sticky="w")
        tk.Label(f3, text="※例: gemini-2.0-flash", font=("Meiryo", 8), bg=bg_color, fg="gray").grid(row=0, column=2, sticky="w")

        self.use_batch_mode = tk.BooleanVar(value=False)
        tk.Checkbutton(f3, text="一括モード（複数ファイルを統合して解析）", variable=self.use_batch_mode, bg=bg_color, font=font_main).grid(row=1, column=0, columnspan=3, sticky="w", pady=2)

        self.use_sqlite_tool = tk.BooleanVar(value=False)
        tk.Checkbutton(f3, text="SQLiteツールを利用", variable=self.use_sqlite_tool, bg=bg_color, font=font_main).grid(row=2, column=0, sticky="w")

        self.use_web_search = tk.BooleanVar(value=False)
        tk.Checkbutton(f3, text="WEB検索を利用", variable=self.use_web_search, bg=bg_color, font=font_main).grid(row=2, column=1, sticky="w")

        self.use_file_search = tk.BooleanVar(value=False)
        tk.Checkbutton(f3, text="RAG(File Search)を利用", variable=self.use_file_search, bg=bg_color, font=font_main).grid(row=3, column=0, sticky="w")

        tk.Label(f3, text="使用RAG:", bg=bg_color, font=font_main).grid(row=3, column=1, sticky="w")
        self.file_search_store = tk.StringVar(value="")
        self.file_search_dropdown = tk.OptionMenu(f3, self.file_search_store, "")
        self.file_search_dropdown.grid(row=4, column=0, columnspan=3, sticky="ew", padx=5, pady=2)
        self.update_rag_list()

        # 4. 実行エリア
        f4 = tk.Frame(self.master, bg=bg_color)
        f4.pack(fill="x", padx=20, pady=10)
        self.run_btn = tk.Button(f4, text="指示ファイルを実行", bg="#ff9999", command=self.run_process, **self.btn_style)
        self.run_btn.pack(side="left")
        tk.Button(f4, text="DBをCSV出力", bg="#99ff99", command=self.export_csv, **self.btn_style).pack(side="left", padx=(20, 0))

        # 5. ログエリア
        log_f = tk.LabelFrame(self.master, text="実行ログ:", bg=bg_color, font=font_main, padx=5, pady=5)
        log_f.pack(fill="both", expand=True, padx=20, pady=5)
        # 修正: 高さを低く設定 (height=6)
        self.log_area = tk.Text(log_f, height=6, font=("Consolas", 9), state="disabled", bg="white")
        self.log_area.pack(side="left", fill="both", expand=True)
        sb = tk.Scrollbar(log_f, command=self.log_area.yview)
        sb.pack(side="right", fill="y")
        self.log_area.config(yscrollcommand=sb.set)

    def log(self, msg):
        now = datetime.datetime.now().strftime("%H:%M:%S")
        self.log_area.config(state="normal")
        self.log_area.insert("end", f"[{now}] {msg}\n")
        self.log_area.see("end")
        self.log_area.config(state="disabled")
        self.master.update()

    def update_rag_list(self):
        if not self.client: return
        try:
            stores = self.client.file_search_stores.list()
            menu = self.file_search_dropdown["menu"]
            menu.delete(0, "end")
            for store in stores:
                store_id = store.name
                menu.add_command(label=store.display_name or store_id, command=lambda x=store_id: self.file_search_store.set(x))
        except: pass

    def select_folder(self):
        path = filedialog.askdirectory()
        if path:
            self.target_folder = path
            self.folder_label.config(text=f"選択フォルダ: {path}")

    def select_prompt_file(self):
        path = filedialog.askopenfilename(filetypes=[("Text Files", "*.txt")])
        if not path: return
        self.prompt_file_path = path
        self.prompt_label.config(text=f"選択ファイル: {path}")
        with open(path, "r", encoding="utf-8") as f:
            self.prompt_text = f.read()
        m = re.search(r"出力項目[:：]\s*(.+)", self.prompt_text)
        if m:
            self.columns = [c.strip() for c in re.split(r'[,、]', m.group(1))]
        else:
            messagebox.showerror("エラー", "プロンプトに『出力項目:』が見つかりません。")

    def show_help(self):
        # 復元: オリジナルの詳細なヘルプテキスト + Pro版の新機能説明
        help_text = (
            "【指示ファイル（プロンプト）の書き方・重要ルール】\n\n"
            "※文字コードは「UTF-8 (BOMなし/UTF-8N)」で保存してください。※\n\n"
            "以下の3つのセクションで構成してください。\n\n"
            "① 概要（任意）\n"
            "   AIに処理対象が何であるかを伝えます。\n\n"
            "② 出力項目（必須・1行で記述）\n"
            "   形式： 出力項目: 項目1, 項目2, 項目3\n"
            "   ※これがデータベースの「列名」になります。\n"
            "   ※項目は「,（半角カンマ）」または「、（全角読点）」で区切ってください。\n\n"
            "③ 処理指示（必須・項目ごとに指定）\n"
            "   各項目を【抽出】するか【生成】するかを明示してください。\n"
            "   ・【抽出】: ファイル内の情報をそのまま抜き出します。検索は行いません。\n"
            "   ・【生成】: AIが推論したり、WEB検索を使って補完したりします。\n\n"
            "【v1.1.0 新機能】\n"
            "・一括モード: 複数ファイル（例：年度別の成績表）を横断して1つのレコードにまとめたい時に使用します。\n"
            "・SQLiteツール: プロンプト内に『###使用するSQLiteのパス』と\"C:\\path\\to\\db.db\"を記述すると、GeminiがDBを参照・更新できます。"
        )
        messagebox.showinfo("TextExtractor1 v1.1.0", help_text)

    def extract_text(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        content = ""
        try:
            if ext in [".txt", ".csv", ".log"]:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f: content = f.read()
            elif ext in [".xlsx", ".xls"]:
                df_dict = pd.read_excel(file_path, sheet_name=None)
                for s, df in df_dict.items(): content += f"\n[Sheet:{s}]\n{df.to_csv(index=False)}"
            elif ext == ".docx":
                doc = docx.Document(file_path)
                content = "\n".join([p.text for p in doc.paragraphs])
            elif ext == ".pptx":
                prs = Presentation(file_path)
                content = "\n".join([shape.text for s in prs.slides for shape in s.shapes if hasattr(shape, "text")])
            elif ext in [".doc", ".ppt", ".jtd"]:
                # 復元: オリジナルの厳密なキーワードリスト
                import olefile
                if olefile.isOleFile(file_path):
                    with olefile.OleFileIO(file_path) as ole:
                        for stream_name_list in ole.listdir():
                            stream_full_path = "/".join(stream_name_list)
                            if any(key in stream_full_path for key in ["WordDocument", "PowerPoint Document", "Current User", "JS-Main", "Body"]):
                                with ole.openstream(stream_name_list) as stream:
                                    data = stream.read()
                                    try: 
                                        decoded_text = data.decode('cp932', errors='ignore')
                                        content += "".join(c for c in decoded_text if c.isprintable() or c in "\n\r\t")
                                    except: 
                                        content += "".join([chr(b) if 32 <= b <= 126 or b >= 128 else " " for b in data])
            elif ext == ".eml":
                with open(file_path, 'rb') as f:
                    msg = BytesParser(policy=policy.default).parse(f)
                content = f"Subject: {msg['subject']}\nFrom: {msg['from']}\nDate: {msg['date']}\n\n"
                content += msg.get_body(preferencelist=('plain')).get_content()
            elif ext == ".msg":
                msg = extract_msg.Message(file_path)
                content = f"Subject: {msg.subject}\nFrom: {msg.sender}\nDate: {msg.date}\n\n{msg.body}"; msg.close()
        except Exception as e:
            self.log(f"テキスト抽出エラー ({ext}): {e}")
        return content

    def sql_tool(self, db_path: str, sql_query: str) -> str:
        try:
            conn = sqlite3.connect(db_path)
            cur = conn.cursor()
            if hasattr(self, 'sqlite_paths'):
                for p in self.sqlite_paths:
                    if os.path.abspath(p) != os.path.abspath(db_path) and os.path.exists(p):
                        cur.execute(f"ATTACH DATABASE '{p}' AS {os.path.splitext(os.path.basename(p))[0]}")
            cur.execute(sql_query)
            if sql_query.strip().upper().startswith("SELECT"):
                res = [dict(zip([d[0] for d in cur.description], r)) for r in cur.fetchall()]
                conn.close()
                return json.dumps(res, ensure_ascii=False)
            else:
                conn.commit()
                res = f"Success: {conn.total_changes} rows affected."
                conn.close()
                return res
        except Exception as e: return f"Error: {str(e)}"

    def get_system_instruction(self):
        # 復元: オリジナルの詳細なシステム指示文
        available_tools = []
        if self.use_web_search.get(): available_tools.append("Google検索")
        if self.use_file_search.get() and self.file_search_store.get(): available_tools.append("RAG（File Search Store）")
        tools_info = "、".join(available_tools) if available_tools else "なし（内蔵知識のみ）"

        instruction = (
            f"\n\n### システムルール (厳守) ###\n"
            f"1. 出力は必ず以下のキーを持つ純粋なJSON形式（またはその配列）のみとしてください。\n"
            f"   キー名: {', '.join(self.columns)}\n"
            f"2. 各項目について、ユーザーの指示にある【抽出】と【生成】の区別を厳密に守ってください。\n"
            f"   ・【抽出】項目: 提示されたファイルの内容のみを使用してください。検索ツールは一切使用せず、見つからない場合は必ず空文字 (\"\") としてください。\n"
            f"   ・【生成】項目: 現在利用可能なツール（{tools_info}）を必要に応じて利用し、あなたの知見や推論に基づいた回答を生成してください。\n"
            f"3. 原則として、英字や記号が多い不自然な文字列はOCR失敗とみなし、日本語として意味が通じるよう再解釈してください。\n"
            f"4. 解説、挨拶、Markdownの装飾(```json等)は一切含めず、パース可能な純粋なJSONデータのみを返してください。\n"
            f"5. 複数の対象がある場合は、必ずJSON配列 [{{...}}, {{...}}] 形式で返してください。\n"
        )
        return instruction

    def run_process(self):
        folder_path = self.folder_label.cget("text").replace("選択フォルダ: ", "")
        if folder_path == "未選択" or not os.path.isdir(folder_path):
            messagebox.showerror("エラー", "有効なフォルダが選択されていません。")
            return
        if not self.prompt_text:
            messagebox.showerror("エラー", "有効なプロンプトファイルが選択されていません。")
            return
        
        self.run_btn.config(state="disabled")
        self.log_area.config(state="normal")
        self.log_area.delete("1.0", "end")
        self.log_area.config(state="disabled")
        self.log("=== 処理開始 ===")

        save_config(self.api_key, self.model_var.get())
        
        self.sqlite_paths = re.findall(r'"([^"]+\.db(?:sqlite\d?)?)"', self.prompt_text)

        db_path = os.path.join(folder_path, f"{os.path.splitext(os.path.basename(self.prompt_file_path))[0]}_DB.db")
        conn = sqlite3.connect(db_path)
        cur = conn.cursor()
        cols_sql = ", ".join([f"'{c}' TEXT" for c in self.columns])
        cur.execute(f"CREATE TABLE IF NOT EXISTS extracted_data (filename TEXT, {cols_sql})")
        conn.commit()

        # 復元: use_recursive に基づくファイルリスト取得ロジック
        files_to_process = []
        if self.use_recursive.get():
            for root_dir, dirs, files in os.walk(folder_path):
                for f in files:
                    if not f.endswith((".db", ".txt")):
                        files_to_process.append(os.path.join(root_dir, f))
        else:
            files_to_process = [os.path.join(folder_path, f) for f in os.listdir(folder_path) 
                                if os.path.isfile(os.path.join(folder_path, f)) and not f.endswith((".db", ".txt"))]
        
        self.log(f"検出項目 ({len(self.columns)}件): {', '.join(self.columns)}")

        try:
            if self.use_batch_mode.get():
                self.process_batch(files_to_process, conn, cur)
            else:
                self.process_individual(files_to_process, conn, cur)
        finally:
            conn.close()
            self.run_btn.config(state="normal")
            self.log("=== 全処理完了 ===")
            messagebox.showinfo("完了", "すべての処理が終了しました。")

    def process_individual(self, files, conn, cur):
        success_count = 0
        fail_count = 0
        for i, file_path in enumerate(files):
            fname = os.path.basename(file_path)
            self.log(f"[{i+1}/{len(files)}] 処理中: {fname}")
            
            cur.execute("SELECT COUNT(*) FROM extracted_data WHERE filename = ?", (fname,))
            if cur.fetchone()[0] > 0:
                self.log("  -> スキップ (処理済み)")
                success_count += 1
                continue

            res = self.call_gemini(file_path)
            if res:
                for item in res:
                    vals = [fname] + [str(item.get(c, "")) for c in self.columns]
                    cur.execute(f"INSERT INTO extracted_data VALUES ({', '.join(['?']*len(vals))})", vals)
                conn.commit()
                self.log(f"  -> 成功: {len(res)}件登録")
                success_count += 1
            else:
                fail_count += 1
            if i < len(files) - 1:
                time.sleep(2)
        self.log(f"結果: 成功 {success_count} / 失敗 {fail_count}")

    def process_batch(self, files, conn, cur):
        CHUNK_SIZE = 20
        accumulated_results = []
        for i in range(0, len(files), CHUNK_SIZE):
            chunk = files[i:i+CHUNK_SIZE]
            self.log(f"一括処理中: {i+1}～{min(i+CHUNK_SIZE, len(files))} ファイル目")
            
            contents = [f"{self.prompt_text}{self.get_system_instruction()}\n"
                        f"※現在【一括モード】で実行中です。以下の複数ファイルを横断して解析してください。\n"]
            
            if accumulated_results:
                contents[0] += f"\n### これまでの解析結果（参考・追記用）:\n{json.dumps(accumulated_results, ensure_ascii=False)}\n"

            parts = []
            for f_path in chunk:
                ext = os.path.splitext(f_path)[1].lower()
                if ext in [".pdf", ".png", ".jpg", ".jpeg", ".webp"]:
                    with open(f_path, "rb") as f:
                        parts.append(types.Part.from_bytes(data=f.read(), mime_type="application/pdf" if ext == ".pdf" else "image/jpeg"))
                else:
                    txt = self.extract_text(f_path)
                    contents[0] += f"\n--- ファイル: {os.path.basename(f_path)} ---\n{txt}\n"

            raw_res = self.call_gemini_raw(contents + parts)
            if raw_res:
                try:
                    parsed = json.loads(re.search(r"(\[.*\]|\{.*\})", raw_res, re.DOTALL).group(1))
                    items = parsed if isinstance(parsed, list) else [parsed]
                    accumulated_results = items
                    # 一括モードでは常に最新の統合結果が正となるため、DBに上書き保存するか運用で決める。
                    # ここではBatch単位の名前で追記していく。
                    for item in items:
                        vals = [f"Batch_{i//CHUNK_SIZE}"] + [str(item.get(c, "")) for c in self.columns]
                        cur.execute(f"INSERT INTO extracted_data VALUES ({', '.join(['?']*len(vals))})", vals)
                    conn.commit()
                    self.log(f"  -> 統合成功")
                except:
                    self.log("  -> 解析失敗 (JSONパースエラー)")
            time.sleep(2)

    def call_gemini(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        instr = f"{self.prompt_text}{self.get_system_instruction()}"
        
        if ext in [".pdf", ".png", ".jpg", ".jpeg", ".webp"]:
            with open(file_path, "rb") as f: fb = f.read()
            mime = "application/pdf" if ext == ".pdf" else "image/jpeg"
            parts = [instr, types.Part.from_bytes(data=fb, mime_type=mime)]
        else:
            txt = self.extract_text(file_path)
            if not txt.strip(): return None
            parts = [f"{instr}\n\n対象内容:\n{txt}"]
        
        raw_res = self.call_gemini_raw(parts)
        if not raw_res: return None
        try:
            parsed = json.loads(re.search(r"(\[.*\]|\{.*\})", raw_res, re.DOTALL).group(1))
            return parsed if isinstance(parsed, list) else [parsed]
        except Exception as e:
            self.log(f"  -> [エラー] 失敗: {e}")
            return None

    def call_gemini_raw(self, contents):
        config = {"response_mime_type": "application/json"}
        if self.use_web_search.get(): config["tools"] = [{"google_search": {}}]
        if self.use_file_search.get() and self.file_search_store.get():
            config.setdefault("tools", []).append({
                "file_search": {"file_search_store_names": [self.file_search_store.get()]}
            })
        if self.use_sqlite_tool.get() and self.sqlite_paths:
            config.setdefault("tools", []).append(self.sql_tool)
        
        # 復元: オリジナルの堅牢なリトライ判定ロジック
        max_retries = 5
        for attempt in range(max_retries):
            try:
                response = self.client.models.generate_content(model=self.model_var.get(), contents=contents, config=config)
                return response.text.strip()
            except Exception as e:
                err_msg = str(e)
                is_retryable = False
                if isinstance(e, APIError) and (e.code in [429, 503, 500, 504]):
                    is_retryable = True
                elif any(x in err_msg for x in ["429", "503", "500", "504", "UNAVAILABLE", "Resource has been exhausted"]):
                    is_retryable = True

                if is_retryable and attempt < max_retries - 1:
                    wait_time = 30 * (2 ** attempt)
                    reason = f"API制限/過負荷"
                    self.log(f"  -> {reason}。{wait_time}秒待機して再試行します... ({attempt + 1}/{max_retries})")
                    time.sleep(wait_time)
                else:
                    self.log(f"  -> [エラー] 失敗: {e}")
                    break
        return None

    def export_csv(self):
        folder_path = self.folder_label.cget("text").replace("選択フォルダ: ", "")
        prompt_file_path = self.prompt_label.cget("text").replace("選択ファイル: ", "")

        if folder_path == "未選択" or prompt_file_path == "未選択":
            messagebox.showerror("エラー", "フォルダと指示ファイルを先に選択してください。")
            return

        prompt_name = os.path.splitext(os.path.basename(prompt_file_path))[0]
        db_path = os.path.join(folder_path, f"{prompt_name}_DB.db")

        if not os.path.exists(db_path):
            messagebox.showerror("エラー", f"対象DBが見つかりません。\nパス: {db_path}\n先に「指示ファイルを実行」してください。")
            return

        default_csv_name = datetime.datetime.now().strftime("%Y%m%d_%H%M%S") + ".csv"
        csv_path = filedialog.asksaveasfilename(
            defaultextension=".csv",
            initialfile=default_csv_name,
            filetypes=[("CSV Files", "*.csv")]
        )
        if not csv_path:
            return

        try:
            conn = sqlite3.connect(db_path)
            cursor = conn.cursor()
            cursor.execute("SELECT * FROM extracted_data")
            rows = cursor.fetchall()
            column_names = [d[0] for d in cursor.description]

            with open(csv_path, "w", encoding="utf_8_sig", newline="") as f:
                writer = csv.writer(f, delimiter=",", quoting=csv.QUOTE_ALL)
                writer.writerow(column_names)
                
                # 復元: セル内改行をスペースに置換する処理
                cleaned_rows = []
                for row in rows:
                    cleaned_rows.append([
                        (v.replace("\r\n", " ").replace("\n", " ").replace("\r", " ") 
                         if isinstance(v, str) else v) 
                        for v in row
                    ])
                writer.writerows(cleaned_rows)

            conn.close()
            messagebox.showinfo("成功", f"CSVファイルを出力しました:\n{csv_path}")
        except Exception as e:
            messagebox.showerror("エラー", f"CSV出力に失敗しました:\n{e}")

if __name__ == "__main__":
    root = tk.Tk()
    app = TextExtractor(root)
    root.mainloop()